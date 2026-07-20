"""Extract text from each slide of a PPTX file and write to markdown."""
from pptx import Presentation
from pptx.util import Inches, Pt
import os

pptx_path = r"d:\DoAn_DaLieu\HUST_PPT_template_2022_RED_16x9_567042.pptx"
output_path = r"d:\DoAn_DaLieu\slide_text.md"

prs = Presentation(pptx_path)

lines = []
for i, slide in enumerate(prs.slides, 1):
    lines.append(f"## Slide {i}:")
    
    slide_texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    slide_texts.append(text)
        
        # Check for tables
        if shape.has_table:
            table = shape.table
            for row in table.rows:
                row_texts = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    row_texts.append(cell_text)
                row_line = " | ".join(row_texts)
                if row_line.strip(" |"):
                    slide_texts.append(f"| {row_line} |")
    
    if slide_texts:
        for t in slide_texts:
            lines.append(t)
    else:
        lines.append("*(Slide trống hoặc chỉ chứa hình ảnh)*")
    
    lines.append("")  # blank line between slides
    lines.append("---")
    lines.append("")

with open(output_path, "w", encoding="utf-8") as f:
    f.write("\n".join(lines))

print(f"Done! Extracted text from {len(prs.slides)} slides to {output_path}")
